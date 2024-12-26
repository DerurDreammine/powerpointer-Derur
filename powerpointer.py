import collections.abc
from pptx import Presentation
from pptx.util import Inches
import json
import re
import random
import time
import prompts

from g4f.client import Client
from g4f.Provider import OpenaiChat, Gemini, DDG
import os
import shutil

os.makedirs("temp", exist_ok=True)

# These are all available designs:
#Design 1 = Envelope, beige
#Design 2 = Blue Bubble
#Design 3 = Light Blue Black
#Design 4 = Black, dark
#Design 5 = wood
#Design 6 = Multicolored, Simple
#Design 7 = Black, white

model_type = "gpt-4o-mini" # gpt-4o-mini gpt-4o gemini-pro

def create_ppt_text(prompt, slides, info=""):
    global model_type
    final_prompt = prompts.make_prompt(prompt, slides, info, model_type)

    # online chat
    client = Client(
      #  api_key="your_api_key_here",
      #  proxies="http://user:pass@host",
    )
    
    # Вопрос и роль                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                        
   # system_role = f"You are an assistant that helps optimize a car (using Arduino Nano) to drive quickly along a black line track. You receive input values: speed=???, maxspeed(on the turn)=???, and sensors_data(every 0.? sec)=?/?/?/?/?/?/?/?/?/?/?/?/?/?/?/?/ where speed is the car's speed on straight sections (1-256), maxspeed is the speed on turns (only the necessary wheel turns at this speed), and sensors_data contains readings from sensors detecting the black line (0 for white, 1 for black). The data format is ?/?/?/?/?/?/?/?/?/?/?/?/?/?/?/?;?/?/?/?/?/?/?/?/?/?/?/?/?/?/?/?, where ';' separates successive readings. For example, 0/0/0/0/0/0/0/1/1/0/0/0/0/0/0/0 indicates driving straight, 0/0/0/0/0/0/0/1/1/1/0/0/0/0/0/0 or 0/0/0/0/0/0/0/1/1/1/1/0/0/0/0/0 indicates turning right, and similar patterns apply for left turns. You control the car using the function move(left_direction(1 for forward, -1 for backward motor rotation), speed, right_direction(1 for forward, -1 for backward motor rotation), speed). Your task is to manage the car's movement based on the input, outputting commands is a only single line in the format '?;move(ld,s,rd,s)/?;move(ld,s,rd,s)/', where '/' is separator, '?' is the duration(delay) in milliseconds, ld and rd are directions (-1 or 1), and s is the motor speed (1..256)! Your primary goal is to navigate the track as quickly as possible by accelerating on long straight sections (add speed to move faster) and slowing down before turns. Remember, you don’t necessarily have to match the sensors_data(every 0.? sec) exactly. If the delays and move commands are identical, you can combine them into a single move with a longer delay. Focus on higher speeds(not inputs maxspeed it is speed 1.256) on straight sections, braking early, entering turns in advance, and optimizing for speed (Avoid using the input speed if acceleration is possible!) (maximal speed is not input maxspeed it is 1..256!)! Reposne only single line in the format '?;move(ld,s,rd,s)/?;move(ld,s,rd,s)/', where '/' is separator, '?' is the duration in milliseconds, ld and rd are directions (-1 or 1), and s is the motor speed (1..256)!" #en
    question = f"{final_prompt}"

    # Формирование сообщения для отправки
    messages = [
      #  {"role": "system", "content": f"{system_role}"},
        {"role": "user", "content": f"{question}"}
    ]

    # Отправка запроса
    response = client.chat.completions.create(
        model=model_type,
        messages=messages,
        provider=DDG
    )

    content = response.choices[0].message.content
    print(content)
    return "Title:" + content

def create_ppt(text_file, design_number, ppt_name):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True
    with open(text_file, 'r', encoding='utf-8') as f: # This is the function for generating the powerpoint. You're a real pro if you understand this lol
        for line_num, line in enumerate(f):
            if line.startswith('Title:'):
                header = line.replace('Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8]
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices)
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue
            elif line.startswith('Header:'):
                header = line.replace('Header:', '').strip()
                continue
            elif line.startswith('Content:'):
                content = line.replace('Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue
                
        if content:
            slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
            title = slide.shapes.title
            title.text = header
            body_shape = slide.shapes.placeholders[slide_placeholder_index]
            tf = body_shape.text_frame
            tf.text = content
            
    os.makedirs("GeneratedPresentations", exist_ok=True)
    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"GeneratedPresentations/{ppt_name}.pptx"
    return f"{file_path}"

def generate_ppt(prompt, add_info, slides, theme):
    prompt = re.sub(r'[^\w\s.\-\(\)]', '', prompt)
    if not theme:
        print("Тема не выбрана, используется стандартная.")
    if theme > 7:
        theme = 1
        print("Неверный номер темы, будет использоватся стандартная.")
    elif theme == 0:
        theme = 1
        print("Неверный номер темы, будет использоватся стандартная.")

    print("Генерация презинтации, это зависит от вашего GPU...\n")
    
    with open(f'temp/{prompt}.txt', 'w', encoding='utf-8') as f:
        f.write(create_ppt_text(prompt, slides, add_info))

    ppt_path = create_ppt(f'temp/{prompt}.txt', theme, prompt)
    return str(ppt_path)

# The main function
def main():
   # print("Welcome to the powerpoint generator!")
 #   topic = input("Topic for the powerpoint: ")
 #   add_info = input("Consider this in the powerpoint (enter if none): ")
    print("Добро пожаловать в генератор презентаций! Улучшено, исправлено и переведено Derur")
    topic = input("Тема презинтации: ")
    add_info = input("Посмотреть это в powerpoint (enter если нет): ") #???????
    if not add_info:
        add_info = ""
   # slides = input("Number of slides: ")
   # theme = int(input("Select theme of the powerpoint (1-7): "))
    slides = input("Количество слайдов: ")
   # theme = int(input("Выберите тему презентации (1-7): "))
    theme = int(input("Выберите стиля презентации (1-7): "))
    start_time = time.time()
  #  print ("Generated and saved under:", generate_ppt(topic, add_info, slides, theme))
    print ("Сгенерировано и сохранено под", generate_ppt(topic, add_info, slides, theme))
    end_time = time.time()
  #  print ("Time used for generating:", round((end_time - start_time), 2))
    print ("Сколько заняла генерация:", round((end_time - start_time), 2))
    if os.path.exists("temp"):
        shutil.rmtree("temp")
    
main()
