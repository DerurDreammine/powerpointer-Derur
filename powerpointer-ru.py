import collections.abc
from pptx import Presentation
from pptx.util import Inches
import json
import re
import random
import time
import prompts

from g4f.client import Client
from g4f.Provider import OpenaiChat, Gemini, DDG
import os
import shutil

os.makedirs("temp", exist_ok=True)

# These are all available designs:
#Design 1 = Envelope, beige
#Design 2 = Blue Bubble
#Design 3 = Light Blue Black
#Design 4 = Black, dark
#Design 5 = wood
#Design 6 = Multicolored, Simple
#Design 7 = Black, white

model_type = "gpt-4o-mini" # gpt-4o-mini gpt-4o gemini-pro

def create_ppt_text(prompt, slides, info=""):
    global model_type
    final_prompt = prompts.make_prompt(prompt, slides, info, model_type)

    # online chat
    client = Client(
      #  api_key="your_api_key_here",
      #  proxies="http://user:pass@host",
    )
    
    # Вопрос и роль                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                                        
   #system_role = f"You are an assistant"
    question = f"{final_prompt}"

    # Формирование сообщения для отправки
    messages = [
      #  {"role": "system", "content": f"{system_role}"},
        {"role": "user", "content": f"{question}"}
    ]

    # Отправка запроса
    response = client.chat.completions.create(
        model=model_type,
        messages=messages,
        provider=DDG
    )

    content = response.choices[0].message.content
    print(content)
    return "Title:" + content

def create_ppt(text_file, design_number, ppt_name):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True
    with open(text_file, 'r', encoding='utf-8') as f: # This is the function for generating the powerpoint. You're a real pro if you understand this lol
        for line_num, line in enumerate(f):
            if line.startswith('Title:'):
                header = line.replace('Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8]
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices)
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue
            elif line.startswith('Header:'):
                header = line.replace('Header:', '').strip()
                continue
            elif line.startswith('Content:'):
                content = line.replace('Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue
                
        if content:
            slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
            title = slide.shapes.title
            title.text = header
            body_shape = slide.shapes.placeholders[slide_placeholder_index]
            tf = body_shape.text_frame
            tf.text = content
            
    os.makedirs("GeneratedPresentations", exist_ok=True)
    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"GeneratedPresentations/{ppt_name}.pptx"
    return f"{file_path}"

def generate_ppt(prompt, add_info, slides, theme):
    prompt = re.sub(r'[^\w\s.\-\(\)]', '', prompt)
    if not theme:
      #  print("No theme selected, using default theme.")
        print("Тема не выбрана, используется стандартная.")
    if theme > 7:
        theme = 1
        #print("Invalid theme number, default theme will be applied.")
        print("Неверный номер темы, будет использоватся стандартная.")
    elif theme == 0:
        theme = 1
   #     print("Invalid theme number, default theme will be applied.")
        print("Неверный номер темы, будет использоватся стандартная.")
    
  #  print("Generating the powerpoint, this could take some time depending on your gpu...\n")
    print("Генерация презинтации, это зависит от вашего GPU...\n")
    
    with open(f'temp/{prompt}.txt', 'w', encoding='utf-8') as f:
        f.write(create_ppt_text(prompt, slides, add_info))

    ppt_path = create_ppt(f'temp/{prompt}.txt', theme, prompt)
    return str(ppt_path)

# The main function
def main():
   # print("Welcome to the powerpoint generator! Updated and fixed Derur")
 #   topic = input("Topic for the powerpoint: ")
 #   add_info = input("Consider this in the powerpoint (enter if none): ")
    print("Добро пожаловать в генератор презентаций! Улучшено, исправлено и переведено Derur")
    topic = input("Тема презинтации: ")
    add_info = input("Посмотреть это в powerpoint (enter если нет): ") #???????
    if not add_info:
        add_info = ""
   # slides = input("Number of slides: ")
   # theme = int(input("Select theme of the powerpoint (1-7): "))
    slides = input("Количество слайдов: ")
   # theme = int(input("Выберите тему презентации (1-7): "))
    theme = int(input("Выберите стиля презентации (1-7): "))
    start_time = time.time()
  #  print ("Generated and saved under:", generate_ppt(topic, add_info, slides, theme))
    print ("Сгенерировано и сохранено под", generate_ppt(topic, add_info, slides, theme))
    end_time = time.time()
  #  print ("Time used for generating:", round((end_time - start_time), 2))
    print ("Сколько заняла генерация:", round((end_time - start_time), 2))
    if os.path.exists("temp"):
        shutil.rmtree("temp")
    
main()
